import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import datetime, timedelta
import os
import openai
from streamlit_calendar import calendar
from dotenv import load_dotenv
import matplotlib.pyplot as plt
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

st.set_page_config(page_title="Resource Monitor", layout="wide")
st.title("📊 Resource Monitoring & GPT Insights App")

uploaded_file = st.sidebar.file_uploader("Upload JIRA Excel", type=["xlsx"])
resources = ["Alice", "Bob", "Charlie", "Diana"]
user = st.sidebar.selectbox("Select Resource", resources)

if "start_time" not in st.session_state:
    st.session_state["start_time"] = datetime.now().time()
if "end_time" not in st.session_state:
    st.session_state["end_time"] = (datetime.now() + timedelta(hours=1)).time()

st.subheader("🛠️ Log Non-Availability")
start_date = st.date_input("Start Date", datetime.today())
start_time = st.time_input("Start Time", st.session_state["start_time"], key="start_time")
end_date = st.date_input("End Date", datetime.today())
end_time = st.time_input("End Time", st.session_state["end_time"], key="end_time")
start_dt = datetime.combine(start_date, start_time)
end_dt = datetime.combine(end_date, end_time)
reason = st.selectbox("Reason", ["Meeting", "Leave", "Sick", "Unplanned Leave", "Out of Office"])

log_file = f"{user}_non_availability.csv"
try:
    log_df = pd.read_csv(log_file)
except FileNotFoundError:
    log_df = pd.DataFrame(columns=["Start", "End", "Reason"])

if st.button("💾 Save Non-Availability"):
    log_df = pd.concat([log_df, pd.DataFrame([{
        "Start": start_dt.strftime("%Y-%m-%d %H:%M:%S"),
        "End": end_dt.strftime("%Y-%m-%d %H:%M:%S"),
        "Reason": reason
    }])], ignore_index=True)
    log_df.to_csv(log_file, index=False)
    st.success("Entry saved!")

calendar_events = []
if not log_df.empty:
    for _, row in log_df.iterrows():
        calendar_events.append({
            "title": f"{row['Reason']}",
            "start": row["Start"],
            "end": row["End"]
        })
calendar(events=calendar_events, options={"editable": False})

csv_summary = ""
summary_data = []
for r in resources:
    path = f"{r}_non_availability.csv"
    if os.path.exists(path):
        df = pd.read_csv(path)
        df["Start"] = pd.to_datetime(df["Start"])
        df["End"] = pd.to_datetime(df["End"])
        df["Hours"] = (df["End"] - df["Start"]).dt.total_seconds() / 3600
        df["Resource"] = r
        summary_data.append(df)
        csv_summary += f"--- {r}_non_availability.csv ---\n" + df.to_string(index=False) + "\n"

if summary_data:
    all_na = pd.concat(summary_data)
    na_summary = all_na.groupby(["Resource", "Reason"]).agg(Total_Hours=("Hours", "sum")).reset_index()
    st.dataframe(na_summary)
    st.bar_chart(na_summary.pivot(index="Resource", columns="Reason", values="Total_Hours").fillna(0))

if uploaded_file:
    df = pd.read_excel(uploaded_file)
    df["Original Estimate (hrs)"] = df["Original Estimate (sec)"] / 3600
    df["Time Spent (hrs)"] = df["Time Spent (sec)"] / 3600
    agg_df = df.groupby("Assignee").agg({
        "Original Estimate (hrs)": "sum",
        "Time Spent (hrs)": "sum"
    }).reset_index()
    agg_df["Utilization (%)"] = (agg_df["Time Spent (hrs)"] / agg_df["Original Estimate (hrs)"]) * 100
    agg_df["Utilization (%)"] = agg_df["Utilization (%)"].round(1)

    availability_map = {}
    for u in df["Assignee"].unique():
        fpath = f"{u}_non_availability.csv"
        if os.path.exists(fpath):
            temp_df = pd.read_csv(fpath)
            temp_df["Start"] = pd.to_datetime(temp_df["Start"])
            temp_df["End"] = pd.to_datetime(temp_df["End"])
            unavailable = (temp_df["End"] - temp_df["Start"]).dt.total_seconds().sum() / 3600
            availability_map[u] = round(80 - unavailable, 2)
        else:
            availability_map[u] = 80

    available_df = pd.DataFrame(list(availability_map.items()), columns=["Assignee", "Available Hours"])
    merged_df = pd.merge(agg_df, available_df, on="Assignee", how="left")
    merged_df["Overallocation Flag"] = merged_df["Original Estimate (hrs)"] > merged_df["Available Hours"]
    st.subheader("📊 Resource Load & Availability")
    st.dataframe(merged_df)

    fig = px.density_heatmap(
        merged_df,
        x="Assignee",
        y="Utilization (%)",
        z="Utilization (%)",
        color_continuous_scale="RdYlGn"
    )
    st.plotly_chart(fig, use_container_width=True)

    # GPT Section
    heatmap_df = merged_df.pivot(index="Assignee", columns="Sprint" if "Sprint" in merged_df.columns else "Available Hours", values="Utilization (%)").fillna(0)
    avg_util = merged_df["Utilization (%)"].mean().round(2)
    max_util = merged_df["Utilization (%)"].max()
    min_util = merged_df["Utilization (%)"].min()
    heatmap_text = heatmap_df.to_string()

    gpt_heatmap_context = f"""
This is a heatmap-like matrix showing resource utilization (%) per assignee per sprint:
{heatmap_text}

Summary:
- Average Utilization: {avg_util}%
- Max Utilization: {max_util}% (likely overloaded)
- Min Utilization: {min_util}% (likely underutilized)
"""

    with st.form("gpt_query_form"):
        query = st.text_area("Ask a question about workload or conflicts:", height=150)
        submitted = st.form_submit_button("🔍 Ask GPT")
        if submitted and query:
            response = openai.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a resource manager assistant."},
                    {"role": "user", "content": gpt_heatmap_context},
                    {"role": "user", "content": query}
                ]
            )
            st.markdown(f"**Answer:** {response.choices[0].message.content}")

    if st.button("📤 Generate PDF Report"):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="Resource Monitoring Report", ln=True, align="C")
        pdf.multi_cell(0, 10, merged_df.to_string(index=False))
        pdf.output("/mnt/data/resource_report.pdf")
        st.success("Report generated!")
        st.download_button("Download PDF", data=open("/mnt/data/resource_report.pdf", "rb"), file_name="resource_report.pdf")

    if st.button("📊 Generate PowerPoint Report"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Resource Utilization Summary"
        temp_fig = plt.figure()
        plt.bar(merged_df["Assignee"], merged_df["Utilization (%)"])
        plt.ylabel("Utilization (%)")
        plt.title("Resource Utilization")
        plt.savefig("/mnt/data/util_chart.png")
        slide.shapes.add_picture("/mnt/data/util_chart.png", Inches(1), Inches(1.5), width=Inches(6))
        prs.save("/mnt/data/resource_report.pptx")
        st.success("PPT report generated.")
        st.download_button("Download PPT", data=open("/mnt/data/resource_report.pptx", "rb"), file_name="resource_report.pptx")